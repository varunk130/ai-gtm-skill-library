"""
GTM Exec Plan - Sample PowerPoint Deck Generator
Produces a polished 4-slide executive GTM deck using Microsoft Fluent design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Microsoft Fluent Design Palette
COLORS = {
    "dark_bg":     RGBColor(0x1B, 0x2A, 0x4A),
    "accent":      RGBColor(0x00, 0x78, 0xD4),
    "accent2":     RGBColor(0x00, 0xB2, 0x94),
    "accent_warm": RGBColor(0xF2, 0x50, 0x22),
    "accent_gold": RGBColor(0xFF, 0xB9, 0x00),
    "text_white":  RGBColor(0xFF, 0xFF, 0xFF),
    "text_dark":   RGBColor(0x2D, 0x2D, 0x2D),
    "text_muted":  RGBColor(0x6B, 0x72, 0x80),
    "light_bg":    RGBColor(0xF3, 0xF4, 0xF6),
    "card_bg":     RGBColor(0xFF, 0xFF, 0xFF),
    "border":      RGBColor(0xE5, 0xE7, 0xEB),
    "purple":      RGBColor(0x5C, 0x2D, 0x91),
    "green":       RGBColor(0x10, 0x7C, 0x10),
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 font_color=COLORS["text_dark"], bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=12, font_color=COLORS["text_dark"],
                  bold=False, space_before=Pt(4), space_after=Pt(2), font_name="Segoe UI"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_metric_card(slide, left, top, width, height, number, label, color):
    card = add_shape(slide, left, top, width, height, COLORS["card_bg"], COLORS["border"])
    # Accent top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # Number
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.55),
                 number, font_size=28, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left + Inches(0.1), top + Inches(0.65), width - Inches(0.2), Inches(0.4),
                 label, font_size=10, font_color=COLORS["text_muted"], alignment=PP_ALIGN.CENTER)


def build_slide_1(prs):
    """Title / Hero Slide - Dark navy background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, COLORS["dark_bg"])

    # Subtle accent line at top
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["accent"]
    line.line.fill.background()

    # Product name
    add_text_box(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(1.2),
                 "M365 Copilot Extensibility", font_size=44, font_color=COLORS["text_white"],
                 bold=True, alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold")

    # Subtitle
    add_text_box(slide, Inches(1.2), Inches(3.0), Inches(9), Inches(0.7),
                 "Go-To-Market Strategy  |  Executive Review", font_size=22,
                 font_color=COLORS["text_muted"], alignment=PP_ALIGN.LEFT)

    # Tagline in accent
    add_text_box(slide, Inches(1.2), Inches(4.0), Inches(9), Inches(0.6),
                 "Build where your users already work", font_size=18,
                 font_color=COLORS["accent2"], bold=False, alignment=PP_ALIGN.LEFT)

    # Bottom bar with metadata
    bottom_bar = add_shape(slide, Inches(0), Inches(6.7), SLIDE_WIDTH, Inches(0.8),
                           RGBColor(0x14, 0x20, 0x3A))
    add_text_box(slide, Inches(1.2), Inches(6.82), Inches(5), Inches(0.5),
                 "April 2026  |  GTM Strategy Team", font_size=11,
                 font_color=COLORS["text_muted"], alignment=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(7), Inches(6.82), Inches(5.5), Inches(0.5),
                 "PRIME Framework  |  AI GTM Skill Library", font_size=11,
                 font_color=COLORS["text_muted"], alignment=PP_ALIGN.RIGHT)

    # Decorative dots
    for i, c in enumerate([COLORS["accent"], COLORS["accent2"], COLORS["purple"], COLORS["accent_gold"]]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2 + i * 0.35), Inches(5.2), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()


def build_slide_2(prs):
    """The Opportunity - Positioning and Readiness"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["light_bg"])

    # Header bar
    header = add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.05), COLORS["dark_bg"])
    add_text_box(slide, Inches(0.8), Inches(0.12), Inches(10), Inches(0.55),
                 "The $12B Extensibility Opportunity Is Ours to Capture",
                 font_size=24, font_color=COLORS["text_white"], bold=True,
                 font_name="Segoe UI Semibold")
    add_text_box(slide, Inches(0.8), Inches(0.62), Inches(8), Inches(0.35),
                 "POSITIONING  AND  READINESS", font_size=10,
                 font_color=COLORS["accent2"], bold=True)

    # Left column - Market context card
    card_l = add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(3.1),
                       COLORS["card_bg"], COLORS["border"])
    add_text_box(slide, Inches(0.7), Inches(1.4), Inches(3), Inches(0.35),
                 "MARKET CONTEXT", font_size=11, font_color=COLORS["accent"], bold=True)

    market_items = [
        ("$47B", "Total enterprise AI platform market (TAM)"),
        ("$12B", "AI extensibility segment addressable (SAM)"),
        ("$1.8B", "Year 1 through M365 commercial base (SOM)"),
    ]
    for i, (num, desc) in enumerate(market_items):
        y = Inches(1.85) + Inches(i * 0.7)
        add_text_box(slide, Inches(0.8), y, Inches(1.2), Inches(0.4),
                     num, font_size=22, font_color=COLORS["accent"], bold=True)
        add_text_box(slide, Inches(2.1), y + Inches(0.05), Inches(3.8), Inches(0.5),
                     desc, font_size=11, font_color=COLORS["text_dark"])

    # Why Now box
    why_card = add_shape(slide, Inches(0.5), Inches(4.55), Inches(5.8), Inches(1.3),
                         COLORS["card_bg"], COLORS["border"])
    add_text_box(slide, Inches(0.7), Inches(4.65), Inches(2), Inches(0.3),
                 "WHY NOW", font_size=11, font_color=COLORS["accent_warm"], bold=True)
    why_items = [
        "68% enterprise AI adoption, customization is #1 unmet need",
        "Competitors 6-12 months behind on enterprise extensibility",
        "M365 Copilot seats create captive distribution channel",
    ]
    for i, item in enumerate(why_items):
        add_text_box(slide, Inches(0.8), Inches(5.0) + Inches(i * 0.28), Inches(5.3), Inches(0.28),
                     f"  {item}", font_size=10, font_color=COLORS["text_dark"])

    # Right column - Positioning card
    card_r = add_shape(slide, Inches(6.6), Inches(1.3), Inches(6.2), Inches(2.0),
                       COLORS["card_bg"], COLORS["border"])
    add_text_box(slide, Inches(6.8), Inches(1.4), Inches(3), Inches(0.35),
                 "POSITIONING", font_size=11, font_color=COLORS["accent"], bold=True)
    add_text_box(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(1.3),
                 "Ship AI agents that reach 400M+ users\nthrough the apps they use every day",
                 font_size=16, font_color=COLORS["dark_bg"], bold=True)

    # Differentiators card
    diff_card = add_shape(slide, Inches(6.6), Inches(3.5), Inches(6.2), Inches(2.35),
                          COLORS["card_bg"], COLORS["border"])
    add_text_box(slide, Inches(6.8), Inches(3.6), Inches(5), Inches(0.3),
                 "KEY DIFFERENTIATORS", font_size=11, font_color=COLORS["accent"], bold=True)

    diffs = [
        ("vs. Company A", "Native M365 Graph, enterprise identity, compliance"),
        ("vs. Company B", "Already where users work: Outlook, Teams, Word"),
        ("vs. Company C", "3x enterprise base, deeper LOB integration"),
        ("vs. Build custom", "10x faster with Copilot Studio + connectors"),
    ]
    for i, (vs, adv) in enumerate(diffs):
        y = Inches(3.95) + Inches(i * 0.45)
        add_text_box(slide, Inches(6.9), y, Inches(2.2), Inches(0.3),
                     vs, font_size=10, font_color=COLORS["purple"], bold=True)
        add_text_box(slide, Inches(9.1), y, Inches(3.5), Inches(0.35),
                     adv, font_size=10, font_color=COLORS["text_dark"])

    # Competitive threat strip at bottom
    strip = add_shape(slide, Inches(0.5), Inches(6.05), Inches(12.3), Inches(1.15),
                      COLORS["dark_bg"])
    add_text_box(slide, Inches(0.8), Inches(6.12), Inches(4), Inches(0.3),
                 "COMPETITIVE LANDSCAPE", font_size=10, font_color=COLORS["accent2"], bold=True)
    comps = [
        ("Company A", "HIGH", COLORS["accent_warm"]),
        ("Company B", "MED", COLORS["accent_gold"]),
        ("Company C", "MED", COLORS["accent_gold"]),
        ("Company D", "MED", COLORS["accent_gold"]),
    ]
    for i, (name, level, color) in enumerate(comps):
        x = Inches(0.8) + Inches(i * 3.1)
        add_text_box(slide, x, Inches(6.42), Inches(2.5), Inches(0.3),
                     name, font_size=11, font_color=COLORS["text_white"], bold=True)
        add_text_box(slide, x, Inches(6.72), Inches(2.5), Inches(0.3),
                     f"Threat: {level}", font_size=9, font_color=color, bold=True)


def build_slide_3(prs):
    """The Plan - Motion and Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["light_bg"])

    # Header
    header = add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.05), COLORS["dark_bg"])
    add_text_box(slide, Inches(0.8), Inches(0.12), Inches(10), Inches(0.55),
                 "500 Extensions and $120M Pipeline in 6 Months",
                 font_size=24, font_color=COLORS["text_white"], bold=True,
                 font_name="Segoe UI Semibold")
    add_text_box(slide, Inches(0.8), Inches(0.62), Inches(8), Inches(0.35),
                 "MOTION  AND  IMPACT", font_size=10,
                 font_color=COLORS["accent2"], bold=True)

    # Metric cards row
    metrics = [
        ("500", "Extensions\n(6-month target)", COLORS["accent"]),
        ("2,000", "Active Developers\n(MAD)", COLORS["accent2"]),
        ("5,000", "Enterprise Orgs\nUsing Extensions", COLORS["purple"]),
        ("$120M", "Partner Co-Sell\nPipeline", COLORS["green"]),
    ]
    for i, (num, label, color) in enumerate(metrics):
        x = Inches(0.5) + Inches(i * 3.15)
        add_metric_card(slide, x, Inches(1.3), Inches(2.9), Inches(1.1), num, label, color)

    # Channel strategy card (left)
    ch_card = add_shape(slide, Inches(0.5), Inches(2.7), Inches(6.0), Inches(3.1),
                        COLORS["card_bg"], COLORS["border"])
    add_text_box(slide, Inches(0.7), Inches(2.8), Inches(4), Inches(0.3),
                 "CHANNEL STRATEGY", font_size=11, font_color=COLORS["accent"], bold=True)

    channels = [
        ("Developer Relations", "40%", "$2.5M", 0.82),
        ("Partner Co-Sell", "30%", "$3.0M", 0.62),
        ("Digital Marketing", "20%", "$1.5M", 0.41),
        ("Events & Conferences", "10%", "$1.0M", 0.20),
    ]
    for i, (ch, pct, inv, bar_w) in enumerate(channels):
        y = Inches(3.2) + Inches(i * 0.6)
        add_text_box(slide, Inches(0.8), y, Inches(2.5), Inches(0.25),
                     ch, font_size=10, font_color=COLORS["text_dark"], bold=True)
        add_text_box(slide, Inches(0.8), y + Inches(0.22), Inches(2), Inches(0.2),
                     f"{pct} contribution  |  {inv}", font_size=9, font_color=COLORS["text_muted"])
        # Bar
        bar_bg = add_shape(slide, Inches(3.5), y + Inches(0.08), Inches(2.8), Inches(0.18),
                           RGBColor(0xE5, 0xE7, 0xEB))
        bar_fg = add_shape(slide, Inches(3.5), y + Inches(0.08), Inches(2.8 * bar_w), Inches(0.18),
                           COLORS["accent"])

    # KPI card (right)
    kpi_card = add_shape(slide, Inches(6.8), Inches(2.7), Inches(5.95), Inches(3.1),
                         COLORS["card_bg"], COLORS["border"])
    add_text_box(slide, Inches(7.0), Inches(2.8), Inches(4), Inches(0.3),
                 "SUCCESS METRICS", font_size=11, font_color=COLORS["accent"], bold=True)

    kpis = [
        ("Monthly Active Developers", "2,000", "by Month 6", "Weekly"),
        ("Extension Publish Rate", "20/week", "steady state", "Weekly"),
        ("Time-to-First-Extension", "< 4 hrs", "developer avg", "Weekly"),
        ("D30 User Retention", "> 40%", "extension users", "Monthly"),
        ("Co-Sell Attach Rate", "15%", "of Copilot deals", "Monthly"),
    ]
    for i, (metric, target, context, cadence) in enumerate(kpis):
        y = Inches(3.2) + Inches(i * 0.48)
        add_text_box(slide, Inches(7.1), y, Inches(2.8), Inches(0.25),
                     metric, font_size=10, font_color=COLORS["text_dark"])
        add_text_box(slide, Inches(10.0), y, Inches(1.2), Inches(0.25),
                     target, font_size=11, font_color=COLORS["accent"], bold=True)
        add_text_box(slide, Inches(11.2), y, Inches(1.4), Inches(0.25),
                     f"{context}  ({cadence})", font_size=8, font_color=COLORS["text_muted"])

    # Partner strip at bottom
    strip = add_shape(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2),
                      COLORS["dark_bg"])
    add_text_box(slide, Inches(0.8), Inches(6.1), Inches(4), Inches(0.25),
                 "PARTNER ECOSYSTEM", font_size=10, font_color=COLORS["accent2"], bold=True)
    partners = [
        ("20 ISV Partners", "Building launch-day extensions"),
        ("5 SI Partners", "Offering implementation services"),
        ("3 Case Studies", "Enterprise early adopter proof points"),
    ]
    for i, (title, desc) in enumerate(partners):
        x = Inches(0.8) + Inches(i * 4.1)
        add_text_box(slide, x, Inches(6.42), Inches(3.5), Inches(0.3),
                     title, font_size=12, font_color=COLORS["text_white"], bold=True)
        add_text_box(slide, x, Inches(6.72), Inches(3.5), Inches(0.3),
                     desc, font_size=9, font_color=COLORS["text_muted"])


def build_slide_4(prs):
    """Execution Roadmap - Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["light_bg"])

    # Header
    header = add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.05), COLORS["dark_bg"])
    add_text_box(slide, Inches(0.8), Inches(0.12), Inches(10), Inches(0.55),
                 "90-Day Launch Sequence with Four Decision Gates",
                 font_size=24, font_color=COLORS["text_white"], bold=True,
                 font_name="Segoe UI Semibold")
    add_text_box(slide, Inches(0.8), Inches(0.62), Inches(8), Inches(0.35),
                 "EXECUTION  ROADMAP", font_size=10,
                 font_color=COLORS["accent2"], bold=True)

    # Timeline connector line
    line = add_shape(slide, Inches(0.9), Inches(2.3), Inches(11.5), Pt(3), COLORS["accent"])

    # Phase cards
    phases = [
        {
            "title": "PRE-LAUNCH",
            "period": "Weeks -4 to -1",
            "color": COLORS["accent"],
            "items": [
                "Finalize messaging + press",
                "Sales enablement package",
                "ISV extensions certified",
                "Developer docs live",
                "Internal readiness training",
            ],
            "gate": "Go / No-Go Review",
        },
        {
            "title": "LAUNCH WEEK",
            "period": "Days 1-5",
            "color": COLORS["accent_warm"],
            "items": [
                "Mon: Press + blog + social",
                "Tue: CEO keynote demo",
                "Wed: Dev livestream",
                "Thu: SI partner launch event",
                "Fri: First metrics review",
            ],
            "gate": "50+ Extensions Published",
        },
        {
            "title": "POST-LAUNCH M1",
            "period": "Days 8-30",
            "color": COLORS["accent2"],
            "items": [
                "Weekly metrics cadence",
                "Community activation",
                "First hackathon",
                "Bug triage + iteration",
                "Partner feedback loop",
            ],
            "gate": "200 MAD, Sentiment > 4.0",
        },
        {
            "title": "SCALE (M2-M3)",
            "period": "Days 31-90",
            "color": COLORS["purple"],
            "items": [
                "Regional expansion",
                "Certification program",
                "Phase 2 investment case",
                "Advanced partner tiers",
                "Competitive response plan",
            ],
            "gate": "1,200 MAD, $40M Pipeline",
        },
    ]

    for i, phase in enumerate(phases):
        x = Inches(0.5) + Inches(i * 3.15)

        # Timeline dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.2), Inches(2.15), Inches(0.35), Inches(0.35))
        dot.fill.solid()
        dot.fill.fore_color.rgb = phase["color"]
        dot.line.fill.background()

        # Phase card
        card = add_shape(slide, x, Inches(2.7), Inches(2.95), Inches(3.2),
                         COLORS["card_bg"], COLORS["border"])

        # Color bar at top of card
        bar = add_shape(slide, x, Inches(2.7), Inches(2.95), Pt(5), phase["color"])

        add_text_box(slide, x + Inches(0.15), Inches(2.82), Inches(2.6), Inches(0.3),
                     phase["title"], font_size=13, font_color=phase["color"], bold=True)
        add_text_box(slide, x + Inches(0.15), Inches(3.1), Inches(2.6), Inches(0.25),
                     phase["period"], font_size=9, font_color=COLORS["text_muted"])

        for j, item in enumerate(phase["items"]):
            add_text_box(slide, x + Inches(0.15), Inches(3.4) + Inches(j * 0.34),
                         Inches(2.6), Inches(0.3),
                         f"  {item}", font_size=9, font_color=COLORS["text_dark"])

        # Gate marker
        gate = add_shape(slide, x + Inches(0.1), Inches(5.2), Inches(2.75), Inches(0.5),
                         phase["color"])
        add_text_box(slide, x + Inches(0.15), Inches(5.22), Inches(2.6), Inches(0.45),
                     f"GATE: {phase['gate']}", font_size=9,
                     font_color=COLORS["text_white"], bold=True, alignment=PP_ALIGN.CENTER)

    # RACI strip at bottom
    strip = add_shape(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.1), COLORS["dark_bg"])
    add_text_box(slide, Inches(0.8), Inches(6.18), Inches(4), Inches(0.25),
                 "OWNERSHIP (RACI)", font_size=10, font_color=COLORS["accent2"], bold=True)

    raci = [
        ("Messaging", "PMM"),
        ("Dev Experience", "DevRel"),
        ("Partner Activation", "Partnerships"),
        ("Sales Enablement", "Enablement"),
        ("Launch Execution", "GTM Lead"),
    ]
    for i, (area, owner) in enumerate(raci):
        x = Inches(0.8) + Inches(i * 2.5)
        add_text_box(slide, x, Inches(6.48), Inches(2.2), Inches(0.25),
                     area, font_size=10, font_color=COLORS["text_white"], bold=True)
        add_text_box(slide, x, Inches(6.73), Inches(2.2), Inches(0.25),
                     f"R: {owner}", font_size=9, font_color=COLORS["text_muted"])


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)

    output_dir = os.path.dirname(os.path.abspath(__file__))
    samples_dir = os.path.join(output_dir, "samples")
    os.makedirs(samples_dir, exist_ok=True)
    output_path = os.path.join(samples_dir, "gtm-plan-m365-copilot-extensibility.pptx")
    prs.save(output_path)
    print(f"Deck saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    main()
